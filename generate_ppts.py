#!/usr/bin/env python3
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

WORKDIR = "/workspace"
ARTICLE_PATH = os.path.join(WORKDIR, "《法華經 藥王菩薩本事品》娑婆堪忍 深體佛意 - 第1809集.txt")
SUMMARY_PATH = os.path.join(WORKDIR, "分析摘要_藥王菩薩本事品_娑婆堪忍_第1809集.txt")
OUT_GENERAL = os.path.join(WORKDIR, "PPT_藥王菩薩本事品_40min.pptx")
OUT_TZUCHI = os.path.join(WORKDIR, "PPT_慈濟資深分享_40min.pptx")


def add_title_slide(prs: Presentation, title: str, subtitle: str, notes: str = ""):
	layout = prs.slide_layouts[0]
	slide = prs.slides.add_slide(layout)
	slide.shapes.title.text = title
	slide.placeholders[1].text = subtitle
	if notes:
		notes_frame = slide.notes_slide.notes_text_frame
		notes_frame.text = notes
	return slide


def add_bullet_slide(prs: Presentation, title: str, bullets: list, notes: str = ""):
	layout = prs.slide_layouts[1]  # Title and Content
	slide = prs.slides.add_slide(layout)
	slide.shapes.title.text = title
	text_frame = slide.shapes.placeholders[1].text_frame
	text_frame.clear()
	for i, bullet in enumerate(bullets):
		p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
		p.text = bullet
		p.level = 0
		for run in p.runs:
			run.font.size = Pt(20)
	if notes:
		notes_frame = slide.notes_slide.notes_text_frame
		notes_frame.text = notes
	return slide


def read_text(path: str) -> str:
	try:
		with open(path, 'r', encoding='utf-8') as f:
			return f.read()
	except Exception as e:
		return f"(讀取失敗: {e})"


def build_general_deck(article_text: str, summary_text: str) -> Presentation:
	prs = Presentation()
	add_title_slide(
		prs,
		"《法華經 藥王菩薩本事品》",
		"娑婆堪忍・深體佛意｜40分鐘分享",
		notes="以〈囑累品〉承先啟後，導入藥王本事之施藥拔苦。強調生活即道場。",
	)

	slides = [
		("議程", [
			"一、背景與問題意識",
			"二、關鍵偈意與主題",
			"三、故事與譬喻",
			"四、佛法義理提要",
			"五、現代觀照與行動",
			"六、結語與落實"
		], "先總覽流程，提醒時間配置：每段約5-7分鐘。"),
		("問題意識：人生八苦", [
			"生、老、病、死與怨憎會等苦皆以身受",
			"苦從觸境而生，觸由執著而起",
			"何處安住？如何離苦？"
		], "連結個人生命經驗，開啟共鳴。"),
		("娑婆堪忍的條件", [
			"此界眾生濁氣重：貪、瞋、癡、慢、疑",
			"五濁惡世：結構性污染與心靈病毒",
			"因緣果報嚴密，需以慧命自度"
		], "強調因果觀與自覺。"),
		("關鍵偈意", [
			"「生生世世為家奴，何時能脫離悲苦？」",
			"小家執著→關係綁縛→輪迴不已",
			"轉出小家，入佛道大家庭"
		], "引用原句引導轉念。"),
		("經典銜接：〈囑累品〉→〈藥王本事〉", [
			"眾會皆大歡喜，未作禮而退",
			"示後有法要：施藥拔苦，轉迷成悟"
		], "點出敘事節點，過渡到行門。"),
		("故事一：多語書生的無依", [
			"才學豐厚不等於生命安住",
			"歷史波瀾與家庭崩解的衝擊",
			"啟示：慧命為依，非智巧可恃"
		], "講述重點：從『能力』到『依止』的轉換。"),
		("故事二：名車與轉念", [
			"富足亦可成『車奴』『宅奴』",
			"轉資源為公義：建院、辦學",
			"啟示：轉念一瞬，福慧雙修"
		], "用對比顯示價值轉換。"),
		("譬喻：煩惱如病毒", [
			"雜惡之毒由心起，蔓延至社會",
			"防疫即防心：斷貪瞋癡的感染鏈"
		], "與疫情經驗對照，易於理解。"),
		("疫情脈絡與因病而貧", [
			"公共衛生危機→經濟脆弱→貧困擴散",
			"慈濟志業：慈善與醫療雙軌並進"
		], "連結現代社會責任。"),
		("佛法義理提要（一）", [
			"生命觀：關係網中的輪迴",
			"因果觀：因緣果報不爽",
			"修行觀：生活即道場"
		], "從三觀建立理解框架。"),
		("佛法義理提要（二）對治七門", [
			"轉念斷執、佈施捨貪、持戒護三業",
			"忍辱化瞋、精進破惰、禪定息散亂",
			"智慧照見如實"
		], "提供可操作的行門線索。"),
		("生活化修行：做中修、修中做", [
			"義診、冬令、環保、護持道場",
			"在日用中鍛鍊堪忍與慈悲"
		], "落地到日常行動。"),
		("離小家入大家", [
			"從自我中心轉為眾生取向",
			"以共同體承載個人長養"
		], "引導加入團體精進。"),
		("操作化建議（一）心法", [
			"每日三觀：無常、因緣、無我",
			"覺察『患得患失』，即起即照"
		], "引導日課養成。"),
		("操作化建議（二）行法", [
			"每週一善行、每月一深修",
			"節律化實踐，積小勝為大勝"
		], "提供節律與節點。"),
		("操作化建議（三）社群", [
			"加入道場、志工團隊",
			"以同行校準與承載"
		], "以人緣成就道業。"),
		("關鍵語彙總結", [
			"堪忍、雜惡、家奴、轉念、施藥",
			"生活即道場、福慧雙修"
		], "重申記憶點。"),
		("Q&A 緩衝", [
			"預留5分鐘互動"
		], "依現場調整時間。"),
		("結語", [
			"娑婆雖苦，願行不退",
			"轉家為道、轉物為法、轉苦為道"
		], "回扣主題與願心。"),
	]

	for title, bullets, notes in slides:
		add_bullet_slide(prs, title, bullets, notes)

	# 附錄：原文摘錄（摘要）
	add_bullet_slide(prs, "附錄：原文摘錄（摘要）", [
		article_text[:200].replace('\n', ' ') + '…',
	], "提醒：更多可參閱全文與分析摘要。")

	# 附錄：分析摘要精要
	summary_preview = [line for line in summary_text.splitlines() if line.strip()][:8]
	add_bullet_slide(prs, "附錄：分析摘要精要", summary_preview, "精要段落供速讀。")

	return prs


def build_tzuchi_deck(article_text: str, summary_text: str) -> Presentation:
	prs = Presentation()
	add_title_slide(
		prs,
		"資深慈濟人的閱讀與實踐",
		"從〈藥王本事〉到志業路｜40分鐘分享",
		notes="站在慈濟志業的節點上，談法入行，行顯法義。",
	)

	slides = [
		("自我定位", [
			"身分：志工/幹部/榮董（自選其一）",
			"角色：承擔、陪伴、傳承",
			"方法：以願領眾、以行導入"
		], "先交代講者在團隊中的責任與視角。"),
		("從經文到實踐的路徑", [
			"觀苦知苦→發悲心→設對治→長願行",
			"以組織保障持續性與擴散性"
		], "把佛法轉化為組織流程。"),
		("四大志業與八大法印", [
			"慈善、醫療、教育、人文",
			"救助、醫療、人本、文化的整合"
		], "連結經義與志業框架。"),
		("案例：冬令發放與義診", [
			"識別需求→動員→執行→回訪",
			"指標：到達率、滿意度、追蹤改善"
		], "呈現流程與數據指標意識。"),
		("案例：環境護持與修樹除草", [
			"做中修，修中做",
			"心靜自涼：以境練心"
		], "引用上人開示語境。"),
		("資源轉化：從物到法", [
			"名車、豪宅→醫院、學校、濟貧",
			"法財兩施，普利群生"
		], "分享捐資轉化的見證。"),
		("志工培力：制度與心法", [
			"訓練：理念、技能、倫理",
			"陪伴：以老帶新，班班相承"
		], "提出可複製的培力模組。"),
		("疫情中的慈濟行動", [
			"醫療量能守護、物資調度、安心關懷",
			"跨域協作與在地韌性"
		], "從『因病而貧』延伸到社會韌性。"),
		("風險與倫理", [
			"個資保護、媒體倫理、資源透明",
			"避免救助依賴，強化自立規劃"
		], "高標準自我要求。"),
		("組織傳承與創新", [
			"制度化×人文化",
			"數位化工具輔助志業管理"
		], "在不失本懷中持續迭代。"),
		("社區連結與跨宗教合作", [
			"共同善行平台",
			"互信互助、擴散善能"
		], "以共善為圓心。"),
		("個人修行節律", [
			"日課：早晚定課、讀書會",
			"月課：閉關日、靜思日"
		], "自我要求與團體節律結合。"),
		("KPI 與成果敘事", [
			"定義：度數、廣度、深度、永續",
			"以人為本的量化與質性並重"
		], "避免唯數字論。"),
		("Q&A 與承諾", [
			"現場互動，凝聚行動承諾"
		], "以願心結尾，呼籲共行。"),
	]

	for title, bullets, notes in slides:
		add_bullet_slide(prs, title, bullets, notes)

	# 附錄：原文摘錄（摘要）
	add_bullet_slide(prs, "附錄：原文摘錄（摘要）", [
		article_text[:200].replace('\n', ' ') + '…',
	], "經文語境作為組織實踐的源頭。")

	# 附錄：分析摘要精要
	summary_preview = [line for line in summary_text.splitlines() if line.strip()][:8]
	add_bullet_slide(prs, "附錄：分析摘要精要", summary_preview, "凝練綱要，便於複誦。")

	return prs


def main():
	article_text = read_text(ARTICLE_PATH)
	summary_text = read_text(SUMMARY_PATH)

	general_deck = build_general_deck(article_text, summary_text)
	general_deck.save(OUT_GENERAL)

	tzuchi_deck = build_tzuchi_deck(article_text, summary_text)
	tzuchi_deck.save(OUT_TZUCHI)

	print(f"已輸出：{OUT_GENERAL}\n已輸出：{OUT_TZUCHI}")


if __name__ == "__main__":
	main()